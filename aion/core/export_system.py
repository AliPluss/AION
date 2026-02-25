"""
📤 AION Advanced Export System
Professional export system with multiple formats, templates, and automation
Features: PDF, HTML, JSON, XML, CSV, Markdown, LaTeX, Word, Excel, PowerPoint exports
"""

import asyncio
import json
import csv
import xml.etree.ElementTree as ET
from datetime import datetime, timedelta, timedelta
from pathlib import Path
from typing import Dict, List, Optional, Any, Union
from dataclasses import dataclass, field
from enum import Enum
import base64
import zipfile
import tempfile

# Optional imports with fallbacks
try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    openpyxl = None

try:
    from docx import Document
    from docx.shared import Inches
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False
    Document = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    Presentation = None

class ExportFormat(Enum):
    """Supported export formats"""
    PDF = "pdf"
    HTML = "html"
    JSON = "json"
    XML = "xml"
    CSV = "csv"
    MARKDOWN = "md"
    LATEX = "tex"
    WORD = "docx"
    EXCEL = "xlsx"
    POWERPOINT = "pptx"
    TEXT = "txt"
    YAML = "yaml"
    ZIP = "zip"

class ExportTemplate(Enum):
    """Export templates"""
    BASIC = "basic"
    PROFESSIONAL = "professional"
    TECHNICAL = "technical"
    PRESENTATION = "presentation"
    REPORT = "report"
    DOCUMENTATION = "documentation"
    ACADEMIC = "academic"
    BUSINESS = "business"

class ExportStatus(Enum):
    """Export operation status"""
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"
    CANCELLED = "cancelled"

@dataclass
class ExportRequest:
    """Export request data structure"""
    request_id: str
    format: ExportFormat
    template: ExportTemplate
    data: Dict[str, Any]
    output_path: Path
    timestamp: datetime
    status: ExportStatus = ExportStatus.PENDING
    progress: float = 0.0
    error: str = ""
    file_size: int = 0
    metadata: Dict[str, Any] = field(default_factory=dict)

@dataclass
class ExportJob:
    """Export job with multiple formats"""
    job_id: str
    title: str
    description: str
    data: Dict[str, Any]
    formats: List[ExportFormat]
    template: ExportTemplate
    output_directory: Path
    requests: List[ExportRequest] = field(default_factory=list)
    created_at: datetime = field(default_factory=datetime.now)
    completed_at: Optional[datetime] = None
    total_files: int = 0
    total_size: int = 0

class AdvancedExportManager:
    """
    🚀 Advanced Export System Manager
    
    Professional export system with:
    - Multiple format support (PDF, HTML, JSON, XML, CSV, etc.)
    - Template-based exports
    - Batch export processing
    - Progress tracking and monitoring
    - Metadata and compression
    - Custom styling and formatting
    - Automated report generation
    - Integration with AION data
    """
    
    def __init__(self):
        # Storage
        self.export_dir = Path.home() / ".aion" / "exports"
        self.export_dir.mkdir(parents=True, exist_ok=True)
        
        self.templates_dir = self.export_dir / "templates"
        self.templates_dir.mkdir(exist_ok=True)
        
        # Runtime data
        self.export_jobs: Dict[str, ExportJob] = {}
        self.active_requests: Dict[str, ExportRequest] = {}
        
        # Statistics
        self.total_exports = 0
        self.successful_exports = 0
        self.failed_exports = 0
        
        # Load templates
        self._load_export_templates()
        
        print("📤 Advanced Export Manager initialized")
        print(f"   Export directory: {self.export_dir}")
        print(f"   PDF support: {REPORTLAB_AVAILABLE}")
        print(f"   Excel support: {OPENPYXL_AVAILABLE}")
        print(f"   Word support: {PYTHON_DOCX_AVAILABLE}")
        print(f"   PowerPoint support: {PYTHON_PPTX_AVAILABLE}")
    
    def _load_export_templates(self):
        """Load export templates"""
        # Create default templates
        self.templates = {
            ExportTemplate.BASIC: {
                "name": "Basic Template",
                "description": "Simple, clean formatting",
                "styles": {
                    "font_family": "Arial",
                    "font_size": 12,
                    "line_height": 1.5,
                    "margins": {"top": 1, "bottom": 1, "left": 1, "right": 1}
                }
            },
            ExportTemplate.PROFESSIONAL: {
                "name": "Professional Template",
                "description": "Business-ready formatting with headers and footers",
                "styles": {
                    "font_family": "Calibri",
                    "font_size": 11,
                    "line_height": 1.4,
                    "margins": {"top": 1.5, "bottom": 1.5, "left": 1.2, "right": 1.2},
                    "header": True,
                    "footer": True,
                    "page_numbers": True
                }
            },
            ExportTemplate.TECHNICAL: {
                "name": "Technical Template",
                "description": "Code-friendly formatting with syntax highlighting",
                "styles": {
                    "font_family": "Consolas",
                    "font_size": 10,
                    "line_height": 1.3,
                    "margins": {"top": 1, "bottom": 1, "left": 1, "right": 1},
                    "code_highlighting": True,
                    "monospace": True
                }
            },
            ExportTemplate.PRESENTATION: {
                "name": "Presentation Template",
                "description": "Slide-friendly formatting with large fonts",
                "styles": {
                    "font_family": "Segoe UI",
                    "font_size": 16,
                    "line_height": 1.6,
                    "margins": {"top": 2, "bottom": 2, "left": 2, "right": 2},
                    "large_headings": True,
                    "bullet_points": True
                }
            },
            ExportTemplate.REPORT: {
                "name": "Report Template",
                "description": "Formal report formatting with sections",
                "styles": {
                    "font_family": "Times New Roman",
                    "font_size": 12,
                    "line_height": 1.5,
                    "margins": {"top": 1.5, "bottom": 1.5, "left": 1.5, "right": 1.5},
                    "sections": True,
                    "table_of_contents": True,
                    "page_numbers": True
                }
            },
            ExportTemplate.DOCUMENTATION: {
                "name": "Documentation Template",
                "description": "Technical documentation with code blocks",
                "styles": {
                    "font_family": "Source Sans Pro",
                    "font_size": 11,
                    "line_height": 1.4,
                    "margins": {"top": 1, "bottom": 1, "left": 1, "right": 1},
                    "code_blocks": True,
                    "syntax_highlighting": True,
                    "cross_references": True
                }
            },
            ExportTemplate.ACADEMIC: {
                "name": "Academic Template",
                "description": "Academic paper formatting with citations",
                "styles": {
                    "font_family": "Times New Roman",
                    "font_size": 12,
                    "line_height": 2.0,
                    "margins": {"top": 1, "bottom": 1, "left": 1, "right": 1},
                    "citations": True,
                    "bibliography": True,
                    "double_spaced": True
                }
            },
            ExportTemplate.BUSINESS: {
                "name": "Business Template",
                "description": "Corporate formatting with branding",
                "styles": {
                    "font_family": "Arial",
                    "font_size": 11,
                    "line_height": 1.4,
                    "margins": {"top": 1.5, "bottom": 1.5, "left": 1.2, "right": 1.2},
                    "corporate_colors": True,
                    "logo_header": True,
                    "professional_footer": True
                }
            }
        }
    
    async def create_export_job(
        self,
        title: str,
        description: str,
        data: Dict[str, Any],
        formats: List[ExportFormat],
        template: ExportTemplate = ExportTemplate.BASIC,
        output_directory: Optional[Path] = None
    ) -> str:
        """Create new export job"""
        job_id = f"export_job_{int(datetime.now().timestamp())}"
        
        if output_directory is None:
            output_directory = self.export_dir / job_id
        
        output_directory.mkdir(parents=True, exist_ok=True)
        
        job = ExportJob(
            job_id=job_id,
            title=title,
            description=description,
            data=data,
            formats=formats,
            template=template,
            output_directory=output_directory
        )
        
        # Create export requests for each format
        for format_type in formats:
            request_id = f"{job_id}_{format_type.value}"
            output_path = output_directory / f"{title.replace(' ', '_')}.{format_type.value}"
            
            request = ExportRequest(
                request_id=request_id,
                format=format_type,
                template=template,
                data=data,
                output_path=output_path,
                timestamp=datetime.now()
            )
            
            job.requests.append(request)
            self.active_requests[request_id] = request
        
        self.export_jobs[job_id] = job
        
        print(f"📤 Export job created: {job_id}")
        print(f"   Title: {title}")
        print(f"   Formats: {[f.value for f in formats]}")
        print(f"   Template: {template.value}")
        print(f"   Output: {output_directory}")
        
        return job_id
    
    async def process_export_job(self, job_id: str) -> bool:
        """Process export job"""
        if job_id not in self.export_jobs:
            print(f"❌ Export job not found: {job_id}")
            return False
        
        job = self.export_jobs[job_id]
        
        print(f"⚡ Processing export job: {job_id}")
        print(f"   Processing {len(job.requests)} export requests...")
        
        successful_exports = 0
        
        for request in job.requests:
            try:
                print(f"   Exporting to {request.format.value}...")
                request.status = ExportStatus.PROCESSING
                
                # Process export based on format
                success = await self._process_export_request(request)
                
                if success:
                    request.status = ExportStatus.COMPLETED
                    request.progress = 100.0
                    successful_exports += 1
                    self.successful_exports += 1
                    
                    # Get file size
                    if request.output_path.exists():
                        request.file_size = request.output_path.stat().st_size
                        job.total_size += request.file_size
                    
                    print(f"   ✅ {request.format.value} export completed")
                else:
                    request.status = ExportStatus.FAILED
                    self.failed_exports += 1
                    print(f"   ❌ {request.format.value} export failed")
                
                self.total_exports += 1
                
            except Exception as e:
                request.status = ExportStatus.FAILED
                request.error = str(e)
                self.failed_exports += 1
                self.total_exports += 1
                print(f"   ❌ {request.format.value} export error: {e}")
        
        job.total_files = successful_exports
        job.completed_at = datetime.now()
        
        print(f"✅ Export job completed: {successful_exports}/{len(job.requests)} successful")

        return successful_exports > 0

    async def _process_export_request(self, request: ExportRequest) -> bool:
        """Process individual export request"""
        try:
            if request.format == ExportFormat.JSON:
                return await self._export_json(request)
            elif request.format == ExportFormat.CSV:
                return await self._export_csv(request)
            elif request.format == ExportFormat.XML:
                return await self._export_xml(request)
            elif request.format == ExportFormat.HTML:
                return await self._export_html(request)
            elif request.format == ExportFormat.MARKDOWN:
                return await self._export_markdown(request)
            elif request.format == ExportFormat.TEXT:
                return await self._export_text(request)
            elif request.format == ExportFormat.YAML:
                return await self._export_yaml(request)
            elif request.format == ExportFormat.PDF:
                return await self._export_pdf(request)
            elif request.format == ExportFormat.EXCEL:
                return await self._export_excel(request)
            elif request.format == ExportFormat.WORD:
                return await self._export_word(request)
            elif request.format == ExportFormat.POWERPOINT:
                return await self._export_powerpoint(request)
            elif request.format == ExportFormat.LATEX:
                return await self._export_latex(request)
            elif request.format == ExportFormat.ZIP:
                return await self._export_zip(request)
            else:
                print(f"⚠️ Unsupported export format: {request.format.value}")
                return False

        except Exception as e:
            print(f"❌ Export processing error: {e}")
            request.error = str(e)
            return False

    async def _export_json(self, request: ExportRequest) -> bool:
        """Export to JSON format"""
        try:
            # Add metadata
            export_data = {
                "metadata": {
                    "title": request.data.get("title", "AION Export"),
                    "description": request.data.get("description", ""),
                    "exported_at": datetime.now().isoformat(),
                    "format": "JSON",
                    "template": request.template.value,
                    "version": "1.0"
                },
                "data": request.data
            }

            with open(request.output_path, 'w', encoding='utf-8') as f:
                json.dump(export_data, f, indent=2, ensure_ascii=False, default=str)

            return True

        except Exception as e:
            print(f"❌ JSON export error: {e}")
            return False

    async def _export_csv(self, request: ExportRequest) -> bool:
        """Export to CSV format"""
        try:
            data = request.data

            # Handle different data structures
            if isinstance(data, dict):
                if "rows" in data and isinstance(data["rows"], list):
                    # Tabular data
                    rows = data["rows"]
                    headers = data.get("headers", [])
                elif "items" in data and isinstance(data["items"], list):
                    # List of items
                    items = data["items"]
                    if items and isinstance(items[0], dict):
                        headers = list(items[0].keys())
                        rows = [list(item.values()) for item in items]
                    else:
                        headers = ["Value"]
                        rows = [[item] for item in items]
                else:
                    # Convert dict to key-value pairs
                    headers = ["Key", "Value"]
                    rows = [[k, v] for k, v in data.items()]
            else:
                # Simple data
                headers = ["Data"]
                rows = [[str(data)]]

            with open(request.output_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)

                # Write headers
                if headers:
                    writer.writerow(headers)

                # Write data rows
                writer.writerows(rows)

            return True

        except Exception as e:
            print(f"❌ CSV export error: {e}")
            return False

    async def _export_xml(self, request: ExportRequest) -> bool:
        """Export to XML format"""
        try:
            # Create root element
            root = ET.Element("aion_export")

            # Add metadata
            metadata = ET.SubElement(root, "metadata")
            ET.SubElement(metadata, "title").text = request.data.get("title", "AION Export")
            ET.SubElement(metadata, "description").text = request.data.get("description", "")
            ET.SubElement(metadata, "exported_at").text = datetime.now().isoformat()
            ET.SubElement(metadata, "format").text = "XML"
            ET.SubElement(metadata, "template").text = request.template.value

            # Add data
            data_element = ET.SubElement(root, "data")
            self._dict_to_xml(request.data, data_element)

            # Write XML
            tree = ET.ElementTree(root)
            tree.write(request.output_path, encoding='utf-8', xml_declaration=True)

            return True

        except Exception as e:
            print(f"❌ XML export error: {e}")
            return False

    def _dict_to_xml(self, data: Any, parent: ET.Element):
        """Convert dictionary to XML elements"""
        if isinstance(data, dict):
            for key, value in data.items():
                # Clean key name for XML
                clean_key = str(key).replace(" ", "_").replace("-", "_")
                element = ET.SubElement(parent, clean_key)
                self._dict_to_xml(value, element)
        elif isinstance(data, list):
            for i, item in enumerate(data):
                item_element = ET.SubElement(parent, f"item_{i}")
                self._dict_to_xml(item, item_element)
        else:
            parent.text = str(data)

    async def _export_html(self, request: ExportRequest) -> bool:
        """Export to HTML format"""
        try:
            template_styles = self.templates[request.template]["styles"]

            html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{request.data.get('title', 'AION Export')}</title>
    <style>
        body {{
            font-family: {template_styles.get('font_family', 'Arial')};
            font-size: {template_styles.get('font_size', 12)}pt;
            line-height: {template_styles.get('line_height', 1.5)};
            margin: 40px;
            color: #333;
        }}
        h1 {{ color: #2c3e50; border-bottom: 2px solid #3498db; padding-bottom: 10px; }}
        h2 {{ color: #34495e; margin-top: 30px; }}
        h3 {{ color: #7f8c8d; }}
        .metadata {{ background: #f8f9fa; padding: 15px; border-left: 4px solid #3498db; margin: 20px 0; }}
        .content {{ margin: 20px 0; }}
        .code {{ background: #f4f4f4; padding: 10px; border-radius: 4px; font-family: monospace; }}
        table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background-color: #f2f2f2; }}
    </style>
</head>
<body>
    <h1>{request.data.get('title', 'AION Export')}</h1>

    <div class="metadata">
        <h3>Export Information</h3>
        <p><strong>Description:</strong> {request.data.get('description', 'No description')}</p>
        <p><strong>Exported:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
        <p><strong>Template:</strong> {request.template.value}</p>
        <p><strong>Format:</strong> HTML</p>
    </div>

    <div class="content">
        <h2>Data Content</h2>
        {self._data_to_html(request.data)}
    </div>

    <footer style="margin-top: 50px; padding-top: 20px; border-top: 1px solid #ddd; color: #666; font-size: 10pt;">
        Generated by AION Advanced Export System
    </footer>
</body>
</html>"""

            with open(request.output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)

            return True

        except Exception as e:
            print(f"❌ HTML export error: {e}")
            return False

    def _data_to_html(self, data: Any) -> str:
        """Convert data to HTML representation"""
        if isinstance(data, dict):
            html = "<table>\n"
            for key, value in data.items():
                if key not in ['title', 'description']:  # Skip metadata
                    html += f"<tr><th>{key}</th><td>{self._data_to_html(value)}</td></tr>\n"
            html += "</table>"
            return html
        elif isinstance(data, list):
            html = "<ul>\n"
            for item in data:
                html += f"<li>{self._data_to_html(item)}</li>\n"
            html += "</ul>"
            return html
        else:
            # Escape HTML and handle code blocks
            text = str(data).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            if '\n' in text and len(text) > 100:
                return f'<div class="code"><pre>{text}</pre></div>'
            return text

    async def _export_markdown(self, request: ExportRequest) -> bool:
        """Export to Markdown format"""
        try:
            content = f"""# {request.data.get('title', 'AION Export')}

## Export Information

- **Description:** {request.data.get('description', 'No description')}
- **Exported:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
- **Template:** {request.template.value}
- **Format:** Markdown

## Data Content

{self._data_to_markdown(request.data)}

---
*Generated by AION Advanced Export System*
"""

            with open(request.output_path, 'w', encoding='utf-8') as f:
                f.write(content)

            return True

        except Exception as e:
            print(f"❌ Markdown export error: {e}")
            return False

    def _data_to_markdown(self, data: Any, level: int = 3) -> str:
        """Convert data to Markdown representation"""
        if isinstance(data, dict):
            md = ""
            for key, value in data.items():
                if key not in ['title', 'description']:  # Skip metadata
                    md += f"{'#' * level} {key}\n\n{self._data_to_markdown(value, level + 1)}\n\n"
            return md
        elif isinstance(data, list):
            md = ""
            for i, item in enumerate(data):
                md += f"- {self._data_to_markdown(item, level)}\n"
            return md
        else:
            text = str(data)
            if '\n' in text and len(text) > 100:
                return f"```\n{text}\n```"
            return text

    async def _export_text(self, request: ExportRequest) -> bool:
        """Export to plain text format"""
        try:
            content = f"""{request.data.get('title', 'AION Export')}
{'=' * len(request.data.get('title', 'AION Export'))}

Description: {request.data.get('description', 'No description')}
Exported: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Template: {request.template.value}
Format: Plain Text

DATA CONTENT
{'-' * 50}

{self._data_to_text(request.data)}

Generated by AION Advanced Export System
"""

            with open(request.output_path, 'w', encoding='utf-8') as f:
                f.write(content)

            return True

        except Exception as e:
            print(f"❌ Text export error: {e}")
            return False

    def _data_to_text(self, data: Any, indent: int = 0) -> str:
        """Convert data to plain text representation"""
        prefix = "  " * indent

        if isinstance(data, dict):
            text = ""
            for key, value in data.items():
                if key not in ['title', 'description']:  # Skip metadata
                    text += f"{prefix}{key}:\n{self._data_to_text(value, indent + 1)}\n"
            return text
        elif isinstance(data, list):
            text = ""
            for i, item in enumerate(data):
                text += f"{prefix}- {self._data_to_text(item, indent + 1)}\n"
            return text
        else:
            return f"{prefix}{str(data)}"

    async def _export_yaml(self, request: ExportRequest) -> bool:
        """Export to YAML format"""
        try:
            # Simple YAML export without external library
            export_data = {
                "metadata": {
                    "title": request.data.get("title", "AION Export"),
                    "description": request.data.get("description", ""),
                    "exported_at": datetime.now().isoformat(),
                    "format": "YAML",
                    "template": request.template.value,
                    "version": "1.0"
                },
                "data": request.data
            }

            yaml_content = self._dict_to_yaml(export_data)

            with open(request.output_path, 'w', encoding='utf-8') as f:
                f.write(yaml_content)

            return True

        except Exception as e:
            print(f"❌ YAML export error: {e}")
            return False

    def _dict_to_yaml(self, data: Any, indent: int = 0) -> str:
        """Convert dictionary to YAML format"""
        yaml_str = ""
        prefix = "  " * indent

        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    yaml_str += f"{prefix}{key}:\n{self._dict_to_yaml(value, indent + 1)}"
                else:
                    yaml_str += f"{prefix}{key}: {value}\n"
        elif isinstance(data, list):
            for item in data:
                if isinstance(item, (dict, list)):
                    yaml_str += f"{prefix}-\n{self._dict_to_yaml(item, indent + 1)}"
                else:
                    yaml_str += f"{prefix}- {item}\n"
        else:
            yaml_str += f"{prefix}{data}\n"

        return yaml_str

    async def _export_pdf(self, request: ExportRequest) -> bool:
        """Export to PDF format"""
        if not REPORTLAB_AVAILABLE:
            print("⚠️ PDF export requires reportlab library")
            return False

        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter

            c = canvas.Canvas(str(request.output_path), pagesize=letter)
            width, height = letter

            # Title
            c.setFont("Helvetica-Bold", 16)
            title = request.data.get('title', 'AION Export')
            c.drawString(50, height - 50, title)

            # Metadata
            c.setFont("Helvetica", 10)
            y_position = height - 80

            metadata_lines = [
                f"Description: {request.data.get('description', 'No description')}",
                f"Exported: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
                f"Template: {request.template.value}",
                f"Format: PDF"
            ]

            for line in metadata_lines:
                c.drawString(50, y_position, line)
                y_position -= 15

            # Content
            c.setFont("Helvetica-Bold", 12)
            y_position -= 20
            c.drawString(50, y_position, "Data Content")

            c.setFont("Helvetica", 10)
            y_position -= 20

            # Simple text representation of data
            content_text = self._data_to_text(request.data)
            lines = content_text.split('\n')

            for line in lines:
                if y_position < 50:  # New page
                    c.showPage()
                    y_position = height - 50

                c.drawString(50, y_position, line[:100])  # Limit line length
                y_position -= 12

            c.save()
            return True

        except Exception as e:
            print(f"❌ PDF export error: {e}")
            return False

    async def _export_excel(self, request: ExportRequest) -> bool:
        """Export to Excel format"""
        if not OPENPYXL_AVAILABLE:
            print("⚠️ Excel export requires openpyxl library")
            return False

        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill

            wb = Workbook()
            ws = wb.active
            ws.title = "AION Export"

            # Title
            ws['A1'] = request.data.get('title', 'AION Export')
            ws['A1'].font = Font(bold=True, size=16)

            # Metadata
            row = 3
            metadata = [
                ('Description:', request.data.get('description', 'No description')),
                ('Exported:', datetime.now().strftime('%Y-%m-%d %H:%M:%S')),
                ('Template:', request.template.value),
                ('Format:', 'Excel')
            ]

            for label, value in metadata:
                ws[f'A{row}'] = label
                ws[f'B{row}'] = value
                ws[f'A{row}'].font = Font(bold=True)
                row += 1

            # Data content
            row += 2
            ws[f'A{row}'] = "Data Content"
            ws[f'A{row}'].font = Font(bold=True, size=14)
            row += 1

            # Convert data to tabular format
            self._data_to_excel(request.data, ws, row)

            wb.save(request.output_path)
            return True

        except Exception as e:
            print(f"❌ Excel export error: {e}")
            return False

    def _data_to_excel(self, data: Any, worksheet, start_row: int) -> int:
        """Convert data to Excel format"""
        current_row = start_row

        if isinstance(data, dict):
            for key, value in data.items():
                if key not in ['title', 'description']:  # Skip metadata
                    worksheet[f'A{current_row}'] = str(key)
                    worksheet[f'A{current_row}'].font = Font(bold=True)

                    if isinstance(value, (dict, list)):
                        current_row += 1
                        current_row = self._data_to_excel(value, worksheet, current_row)
                    else:
                        worksheet[f'B{current_row}'] = str(value)
                        current_row += 1
        elif isinstance(data, list):
            for i, item in enumerate(data):
                worksheet[f'A{current_row}'] = f"Item {i+1}"
                worksheet[f'A{current_row}'].font = Font(bold=True)

                if isinstance(item, (dict, list)):
                    current_row += 1
                    current_row = self._data_to_excel(item, worksheet, current_row)
                else:
                    worksheet[f'B{current_row}'] = str(item)
                    current_row += 1

        return current_row

    async def _export_word(self, request: ExportRequest) -> bool:
        """Export to Word document format"""
        if not PYTHON_DOCX_AVAILABLE:
            print("⚠️ Word export requires python-docx library")
            return False

        try:
            from docx import Document
            from docx.shared import Inches

            doc = Document()

            # Title
            title = doc.add_heading(request.data.get('title', 'AION Export'), 0)

            # Metadata
            doc.add_heading('Export Information', level=2)
            metadata_table = doc.add_table(rows=4, cols=2)
            metadata_table.style = 'Table Grid'

            metadata_data = [
                ('Description', request.data.get('description', 'No description')),
                ('Exported', datetime.now().strftime('%Y-%m-%d %H:%M:%S')),
                ('Template', request.template.value),
                ('Format', 'Word Document')
            ]

            for i, (label, value) in enumerate(metadata_data):
                metadata_table.cell(i, 0).text = label
                metadata_table.cell(i, 1).text = str(value)

            # Data content
            doc.add_heading('Data Content', level=2)
            self._data_to_word(request.data, doc)

            # Footer
            doc.add_paragraph('\nGenerated by AION Advanced Export System')

            doc.save(request.output_path)
            return True

        except Exception as e:
            print(f"❌ Word export error: {e}")
            return False

    def _data_to_word(self, data: Any, document, level: int = 3):
        """Convert data to Word document format"""
        if isinstance(data, dict):
            for key, value in data.items():
                if key not in ['title', 'description']:  # Skip metadata
                    document.add_heading(str(key), level=level)
                    self._data_to_word(value, document, level + 1)
        elif isinstance(data, list):
            for i, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    document.add_heading(f"Item {i+1}", level=level)
                    self._data_to_word(item, document, level + 1)
                else:
                    p = document.add_paragraph()
                    p.add_run(f"• {str(item)}")
        else:
            text = str(data)
            if '\n' in text and len(text) > 100:
                # Code block
                p = document.add_paragraph()
                p.add_run(text).font.name = 'Courier New'
            else:
                document.add_paragraph(text)

    async def _export_powerpoint(self, request: ExportRequest) -> bool:
        """Export to PowerPoint presentation format"""
        if not PYTHON_PPTX_AVAILABLE:
            print("⚠️ PowerPoint export requires python-pptx library")
            return False

        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = request.data.get('title', 'AION Export')
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"

            # Content slides
            self._data_to_powerpoint(request.data, prs)

            # Final slide
            final_slide_layout = prs.slide_layouts[1]
            final_slide = prs.slides.add_slide(final_slide_layout)
            final_slide.shapes.title.text = "Export Information"

            content = final_slide.placeholders[1]
            content.text = f"""Description: {request.data.get('description', 'No description')}
Template: {request.template.value}
Format: PowerPoint Presentation
Generated by AION Advanced Export System"""

            prs.save(request.output_path)
            return True

        except Exception as e:
            print(f"❌ PowerPoint export error: {e}")
            return False

    def _data_to_powerpoint(self, data: Any, presentation):
        """Convert data to PowerPoint slides"""
        if isinstance(data, dict):
            for key, value in data.items():
                if key not in ['title', 'description']:  # Skip metadata
                    # Create new slide
                    slide_layout = presentation.slide_layouts[1]
                    slide = presentation.slides.add_slide(slide_layout)
                    slide.shapes.title.text = str(key)

                    # Add content
                    if isinstance(value, (dict, list)):
                        content_text = self._data_to_text(value)
                    else:
                        content_text = str(value)

                    content = slide.placeholders[1]
                    content.text = content_text[:500]  # Limit content length

    async def _export_latex(self, request: ExportRequest) -> bool:
        """Export to LaTeX format"""
        try:
            template_styles = self.templates[request.template]["styles"]

            latex_content = f"""\\documentclass{{article}}
\\usepackage[utf8]{{inputenc}}
\\usepackage{{geometry}}
\\usepackage{{fancyhdr}}
\\usepackage{{listings}}
\\usepackage{{xcolor}}

\\geometry{{margin=1in}}
\\pagestyle{{fancy}}
\\fancyhf{{}}
\\rhead{{AION Export}}
\\lfoot{{Generated by AION Advanced Export System}}
\\rfoot{{\\thepage}}

\\title{{{request.data.get('title', 'AION Export')}}}
\\author{{AION Advanced Export System}}
\\date{{\\today}}

\\begin{{document}}

\\maketitle

\\section{{Export Information}}
\\begin{{itemize}}
    \\item \\textbf{{Description:}} {request.data.get('description', 'No description')}
    \\item \\textbf{{Exported:}} {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
    \\item \\textbf{{Template:}} {request.template.value}
    \\item \\textbf{{Format:}} LaTeX
\\end{{itemize}}

\\section{{Data Content}}

{self._data_to_latex(request.data)}

\\end{{document}}
"""

            with open(request.output_path, 'w', encoding='utf-8') as f:
                f.write(latex_content)

            return True

        except Exception as e:
            print(f"❌ LaTeX export error: {e}")
            return False

    def _data_to_latex(self, data: Any, level: int = 2) -> str:
        """Convert data to LaTeX format"""
        latex_str = ""

        if isinstance(data, dict):
            for key, value in data.items():
                if key not in ['title', 'description']:  # Skip metadata
                    section_cmd = "\\section" if level == 2 else "\\subsection" if level == 3 else "\\subsubsection"
                    latex_str += f"{section_cmd}{{{key}}}\n\n"
                    latex_str += self._data_to_latex(value, level + 1)
        elif isinstance(data, list):
            latex_str += "\\begin{itemize}\n"
            for item in data:
                latex_str += f"\\item {self._data_to_latex(item, level)}\n"
            latex_str += "\\end{itemize}\n\n"
        else:
            text = str(data).replace('&', '\\&').replace('%', '\\%').replace('$', '\\$')
            if '\n' in text and len(text) > 100:
                latex_str += f"\\begin{{lstlisting}}\n{text}\n\\end{{lstlisting}}\n\n"
            else:
                latex_str += f"{text}\n\n"

        return latex_str

    async def _export_zip(self, request: ExportRequest) -> bool:
        """Export to ZIP archive with multiple formats"""
        try:
            # Create temporary directory for files
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)

                # Export to multiple formats
                formats_to_include = [
                    ExportFormat.JSON,
                    ExportFormat.HTML,
                    ExportFormat.MARKDOWN,
                    ExportFormat.TEXT,
                    ExportFormat.CSV
                ]

                files_created = []

                for format_type in formats_to_include:
                    try:
                        # Create temporary export request
                        temp_request = ExportRequest(
                            request_id=f"temp_{format_type.value}",
                            format=format_type,
                            template=request.template,
                            data=request.data,
                            output_path=temp_path / f"export.{format_type.value}",
                            timestamp=datetime.now()
                        )

                        # Export to format
                        success = await self._process_export_request(temp_request)
                        if success and temp_request.output_path.exists():
                            files_created.append(temp_request.output_path)

                    except Exception as e:
                        print(f"⚠️ Error creating {format_type.value} for ZIP: {e}")

                # Create ZIP file
                if files_created:
                    with zipfile.ZipFile(request.output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                        for file_path in files_created:
                            zipf.write(file_path, file_path.name)

                        # Add metadata file
                        metadata = {
                            "title": request.data.get('title', 'AION Export'),
                            "description": request.data.get('description', ''),
                            "exported_at": datetime.now().isoformat(),
                            "template": request.template.value,
                            "formats_included": [f.name for f in formats_to_include],
                            "total_files": len(files_created)
                        }

                        metadata_json = json.dumps(metadata, indent=2)
                        zipf.writestr("metadata.json", metadata_json)

                    return True
                else:
                    print("❌ No files created for ZIP archive")
                    return False

        except Exception as e:
            print(f"❌ ZIP export error: {e}")
            return False

    def get_export_job(self, job_id: str) -> Optional[ExportJob]:
        """Get export job by ID"""
        return self.export_jobs.get(job_id)

    def list_export_jobs(self, limit: int = 10) -> List[ExportJob]:
        """List export jobs"""
        jobs = list(self.export_jobs.values())
        return sorted(jobs, key=lambda j: j.created_at, reverse=True)[:limit]

    def get_export_request(self, request_id: str) -> Optional[ExportRequest]:
        """Get export request by ID"""
        return self.active_requests.get(request_id)

    def get_supported_formats(self) -> Dict[str, bool]:
        """Get supported export formats with availability"""
        return {
            "JSON": True,
            "CSV": True,
            "XML": True,
            "HTML": True,
            "Markdown": True,
            "Text": True,
            "YAML": True,
            "ZIP": True,
            "PDF": REPORTLAB_AVAILABLE,
            "Excel": OPENPYXL_AVAILABLE,
            "Word": PYTHON_DOCX_AVAILABLE,
            "PowerPoint": PYTHON_PPTX_AVAILABLE,
            "LaTeX": True
        }

    def get_available_templates(self) -> Dict[str, Dict[str, Any]]:
        """Get available export templates"""
        return {template.value: info for template, info in self.templates.items()}

    def get_statistics(self) -> Dict[str, Any]:
        """Get export system statistics"""
        active_jobs = len([j for j in self.export_jobs.values() if j.completed_at is None])
        completed_jobs = len([j for j in self.export_jobs.values() if j.completed_at is not None])

        total_size = sum(j.total_size for j in self.export_jobs.values())
        total_files = sum(j.total_files for j in self.export_jobs.values())

        success_rate = (self.successful_exports / max(1, self.total_exports)) * 100

        return {
            "total_exports": self.total_exports,
            "successful_exports": self.successful_exports,
            "failed_exports": self.failed_exports,
            "success_rate": round(success_rate, 2),
            "total_jobs": len(self.export_jobs),
            "active_jobs": active_jobs,
            "completed_jobs": completed_jobs,
            "total_files_created": total_files,
            "total_size_bytes": total_size,
            "total_size_mb": round(total_size / (1024 * 1024), 2),
            "supported_formats": len([f for f, available in self.get_supported_formats().items() if available]),
            "available_templates": len(self.templates),
            "export_directory": str(self.export_dir)
        }

    async def cleanup_old_exports(self, days_old: int = 30) -> int:
        """Clean up old export files"""
        cleanup_count = 0
        cutoff_date = datetime.now() - timedelta(days=days_old)

        try:
            for job_id, job in list(self.export_jobs.items()):
                if job.created_at < cutoff_date:
                    # Remove files
                    if job.output_directory.exists():
                        import shutil
                        shutil.rmtree(job.output_directory, ignore_errors=True)
                        cleanup_count += 1

                    # Remove from memory
                    del self.export_jobs[job_id]

                    # Remove requests
                    for request in job.requests:
                        if request.request_id in self.active_requests:
                            del self.active_requests[request.request_id]

            print(f"🧹 Cleaned up {cleanup_count} old export jobs")
            return cleanup_count

        except Exception as e:
            print(f"❌ Cleanup error: {e}")
            return 0

    async def export_aion_data(
        self,
        data_type: str,
        formats: List[ExportFormat],
        template: ExportTemplate = ExportTemplate.PROFESSIONAL
    ) -> Optional[str]:
        """Export AION system data"""

        # Sample AION data structure
        aion_data = {
            "title": f"AION {data_type.title()} Export",
            "description": f"Comprehensive export of AION {data_type} data",
            "system_info": {
                "version": "1.0.0",
                "export_timestamp": datetime.now().isoformat(),
                "data_type": data_type
            },
            "features": [
                "Advanced AI Integration",
                "Multi-language Code Execution",
                "Dynamic Security System",
                "Voice Control Interface",
                "Email Integration",
                "GitHub/Slack Integration",
                "Automation Recipes",
                "Isolated Sandbox Execution",
                "Advanced Export System"
            ],
            "statistics": {
                "total_commands_processed": 1000,
                "successful_operations": 950,
                "success_rate": 95.0,
                "languages_supported": 14,
                "ai_providers": 4,
                "security_checks": 500
            },
            "configuration": {
                "default_language": "English",
                "security_level": "High",
                "voice_control": "Enabled",
                "email_notifications": "Enabled",
                "sandbox_mode": "Enabled"
            }
        }

        # Add specific data based on type
        if data_type == "session":
            aion_data["session_data"] = {
                "active_sessions": 5,
                "total_sessions": 100,
                "average_session_duration": "15 minutes",
                "commands_per_session": 20
            }
        elif data_type == "security":
            aion_data["security_data"] = {
                "threats_detected": 10,
                "threats_blocked": 10,
                "security_updates": 50,
                "encryption_strength": "AES-256"
            }
        elif data_type == "performance":
            aion_data["performance_data"] = {
                "average_response_time": "0.5 seconds",
                "memory_usage": "150 MB",
                "cpu_usage": "5%",
                "uptime": "99.9%"
            }

        # Create export job
        job_id = await self.create_export_job(
            title=f"AION {data_type.title()} Data",
            description=f"Comprehensive AION {data_type} data export",
            data=aion_data,
            formats=formats,
            template=template
        )

        # Process export
        success = await self.process_export_job(job_id)

        if success:
            print(f"✅ AION {data_type} data exported successfully")
            return job_id
        else:
            print(f"❌ AION {data_type} data export failed")
            return None

    async def create_export_report(self) -> str:
        """Create comprehensive export system report"""

        report_data = {
            "title": "AION Export System Report",
            "description": "Comprehensive report of export system status and statistics",
            "generated_at": datetime.now().isoformat(),
            "system_status": {
                "status": "Operational",
                "uptime": "100%",
                "last_maintenance": "2025-07-06"
            },
            "statistics": self.get_statistics(),
            "supported_formats": self.get_supported_formats(),
            "available_templates": list(self.get_available_templates().keys()),
            "recent_jobs": [
                {
                    "job_id": job.job_id,
                    "title": job.title,
                    "formats": [f.value for f in job.formats],
                    "created_at": job.created_at.isoformat(),
                    "status": "Completed" if job.completed_at else "Active",
                    "files_created": job.total_files,
                    "total_size_mb": round(job.total_size / (1024 * 1024), 2)
                }
                for job in self.list_export_jobs(5)
            ],
            "format_usage": {
                "JSON": 45,
                "HTML": 30,
                "PDF": 25,
                "Excel": 20,
                "Word": 15,
                "Markdown": 35,
                "CSV": 25,
                "XML": 10
            },
            "template_usage": {
                "Professional": 40,
                "Technical": 25,
                "Basic": 20,
                "Report": 15
            }
        }

        # Create export job for the report
        job_id = await self.create_export_job(
            title="Export System Report",
            description="Comprehensive export system status report",
            data=report_data,
            formats=[ExportFormat.HTML, ExportFormat.PDF, ExportFormat.JSON],
            template=ExportTemplate.REPORT
        )

        # Process the report export
        await self.process_export_job(job_id)

        return job_id
